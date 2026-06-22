from pptx import Presentation
from pptx.util import Inches, Pt


prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
slide = prs.slides.add_slide(prs.slide_layouts[6])
box = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(5), Inches(1))
run = box.text_frame.paragraphs[0].add_run()
run.text = "CyberLab PPT diagram test"
run.font.size = Pt(28)
prs.save("artifacts/generated/diagrams/test_diagram.pptx")
