from pathlib import Path

from PIL import Image, ImageDraw, ImageFont
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_CONNECTOR, MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt


OUT = Path("artifacts/generated/diagrams")
PPTX_DIR = OUT / "pptx"
PPTX_DIR.mkdir(parents=True, exist_ok=True)

W = Inches(13.333)
H = Inches(7.5)

NAVY = RGBColor(15, 42, 67)
BLUE = RGBColor(36, 99, 170)
TEAL = RGBColor(15, 118, 110)
GOLD = RGBColor(180, 120, 30)
PURPLE = RGBColor(103, 80, 164)
RED = RGBColor(155, 28, 28)
INK = RGBColor(31, 41, 55)
MUTED = RGBColor(100, 116, 139)
BORDER = RGBColor(203, 213, 225)
FILL = RGBColor(248, 250, 252)
BLUE_FILL = RGBColor(232, 241, 252)
TEAL_FILL = RGBColor(228, 247, 244)
GOLD_FILL = RGBColor(255, 247, 225)
PURPLE_FILL = RGBColor(243, 238, 255)
RED_FILL = RGBColor(255, 241, 242)

PNG_W = 1920
PNG_H = 1084


def pil_color(color):
    return (color[0], color[1], color[2])


def cn_font(size, bold=False):
    candidates = [
        "/System/Library/Fonts/PingFang.ttc",
        "/System/Library/Fonts/STHeiti Light.ttc",
        "/Library/Fonts/Arial Unicode.ttf",
    ]
    for path in candidates:
        if Path(path).exists():
            try:
                return ImageFont.truetype(path, size=size)
            except Exception:
                pass
    return ImageFont.load_default()


def draw_center_text(draw, box, text, font, fill=(31, 41, 55), line_gap=6):
    x1, y1, x2, y2 = box
    lines = text.split("\n")
    heights = []
    widths = []
    for line in lines:
        bbox = draw.textbbox((0, 0), line, font=font)
        widths.append(bbox[2] - bbox[0])
        heights.append(bbox[3] - bbox[1])
    total_h = sum(heights) + line_gap * (len(lines) - 1)
    y = y1 + (y2 - y1 - total_h) / 2
    for line, w, h in zip(lines, widths, heights):
        x = x1 + (x2 - x1 - w) / 2
        draw.text((x, y), line, font=font, fill=fill)
        y += h + line_gap


def png_rect(draw, box, text, fill, outline, text_fill=(31, 41, 55), size=28, bold=False, radius=22, width=3):
    draw.rounded_rectangle(box, radius=radius, fill=fill, outline=outline, width=width)
    draw_center_text(draw, box, text, cn_font(size, bold), text_fill)


def png_arrow(draw, start, end, fill=(36, 99, 170), width=5):
    draw.line([start, end], fill=fill, width=width)
    x1, y1 = start
    x2, y2 = end
    if x2 >= x1:
        points = [(x2, y2), (x2 - 20, y2 - 12), (x2 - 20, y2 + 12)]
    else:
        points = [(x2, y2), (x2 + 20, y2 - 12), (x2 + 20, y2 + 12)]
    draw.polygon(points, fill=fill)


def png_title(draw, title_text, subtitle_text):
    draw.text((78, 56), title_text, font=cn_font(42, True), fill=pil_color((15, 42, 67)))
    draw.text((82, 118), subtitle_text, font=cn_font(22), fill=pil_color((100, 116, 139)))
    draw.rectangle((78, 166, 1840, 171), fill=pil_color((36, 99, 170)))


def new_prs():
    prs = Presentation()
    prs.slide_width = W
    prs.slide_height = H
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg = slide.background.fill
    bg.solid()
    bg.fore_color.rgb = RGBColor(255, 255, 255)
    return prs, slide


def set_text(shape, text, size=16, bold=False, color=INK, align=PP_ALIGN.CENTER):
    tf = shape.text_frame
    tf.clear()
    tf.margin_left = Inches(0.08)
    tf.margin_right = Inches(0.08)
    tf.margin_top = Inches(0.04)
    tf.margin_bottom = Inches(0.04)
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.name = "PingFang SC"
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color


def title(slide, text, subtitle=None):
    box = slide.shapes.add_textbox(Inches(0.55), Inches(0.28), Inches(9.2), Inches(0.48))
    set_text(box, text, size=24, bold=True, color=NAVY, align=PP_ALIGN.LEFT)
    if subtitle:
        sub = slide.shapes.add_textbox(Inches(0.58), Inches(0.82), Inches(10.5), Inches(0.32))
        set_text(sub, subtitle, size=10.5, color=MUTED, align=PP_ALIGN.LEFT)
    rule = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.55), Inches(1.15), Inches(12.2), Inches(0.03))
    rule.fill.solid()
    rule.fill.fore_color.rgb = BLUE
    rule.line.fill.background()


def footer(slide):
    box = slide.shapes.add_textbox(Inches(10.2), Inches(7.05), Inches(2.5), Inches(0.22))
    set_text(box, "CyberLab ContentOps 详细设计", size=8.5, color=MUTED, align=PP_ALIGN.RIGHT)


def rect(slide, x, y, w, h, text, fill=FILL, line=BORDER, color=INK, size=13, bold=False):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    shape.line.color.rgb = line
    shape.line.width = Pt(1.2)
    set_text(shape, text, size=size, bold=bold, color=color)
    return shape


def label(slide, x, y, w, h, text, color=MUTED, size=10, align=PP_ALIGN.LEFT):
    shape = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    set_text(shape, text, size=size, color=color, align=align)
    return shape


def arrow(slide, x1, y1, x2, y2, color=BLUE, width=1.6):
    conn = slide.shapes.add_connector(
        MSO_CONNECTOR.STRAIGHT,
        Inches(x1),
        Inches(y1),
        Inches(x2),
        Inches(y2),
    )
    conn.line.color.rgb = color
    conn.line.width = Pt(width)
    try:
        conn.line.end_arrowhead = True
    except Exception:
        pass
    return conn


def save(prs, name):
    prs.save(PPTX_DIR / f"{name}.pptx")


def system_architecture():
    prs, s = new_prs()
    title(s, "系统总体调用图", "前端只调用后端 API，外部平台工具统一收束到后端适配器边界")
    footer(s)

    rect(s, 0.75, 1.65, 2.25, 0.75, "React/Vite\n前端控制台", BLUE_FILL, BLUE, BLUE, bold=True)
    rect(s, 3.45, 1.65, 2.25, 0.75, "FastAPI\nAPI 路由", BLUE_FILL, BLUE, BLUE, bold=True)
    rect(s, 6.15, 1.25, 2.2, 0.62, "TrendService\n趋势采集", TEAL_FILL, TEAL, TEAL, size=12, bold=True)
    rect(s, 6.15, 2.05, 2.2, 0.62, "ContentService\n内容生成", TEAL_FILL, TEAL, TEAL, size=12, bold=True)
    rect(s, 6.15, 2.85, 2.2, 0.62, "AdapterRegistry\n平台适配", TEAL_FILL, TEAL, TEAL, size=12, bold=True)
    rect(s, 6.15, 3.65, 2.2, 0.62, "MemoryRepository\n运行状态", TEAL_FILL, TEAL, TEAL, size=12, bold=True)

    rect(s, 9.05, 1.25, 2.95, 0.62, "wemp-operator\n热点采集脚本", GOLD_FILL, GOLD, GOLD, size=12, bold=True)
    rect(s, 9.05, 2.05, 2.95, 0.62, "XiaohongshuSkills\nCDP 预览边界", PURPLE_FILL, PURPLE, PURPLE, size=12, bold=True)
    rect(s, 9.05, 2.85, 2.95, 0.62, "toutiao CLI/MCP\n今日头条预览边界", PURPLE_FILL, PURPLE, PURPLE, size=12, bold=True)
    rect(s, 9.05, 3.65, 2.95, 0.62, "微信公众号\nMarkdown/草稿边界", PURPLE_FILL, PURPLE, PURPLE, size=12, bold=True)

    arrow(s, 3.0, 2.02, 3.45, 2.02)
    arrow(s, 5.7, 2.02, 6.15, 1.56)
    arrow(s, 5.7, 2.02, 6.15, 2.36)
    arrow(s, 5.7, 2.02, 6.15, 3.16)
    arrow(s, 5.7, 2.02, 6.15, 3.96)
    arrow(s, 8.35, 1.56, 9.05, 1.56, TEAL)
    arrow(s, 8.35, 3.16, 9.05, 2.36, TEAL)
    arrow(s, 8.35, 3.16, 9.05, 3.16, TEAL)
    arrow(s, 8.35, 3.16, 9.05, 3.96, TEAL)

    rect(s, 0.78, 5.05, 11.25, 0.95, "安全边界：当前仅生成 preview / draft 载荷，不自动点击真实平台发布按钮；真实动作必须经过人工审核门禁。", RED_FILL, RED, RED, size=13, bold=True)
    save(prs, "01_system_architecture")


def output_flow():
    prs, s = new_prs()
    title(s, "详细设计输出流", "从输入参数到实验报告，每一步都保留可审计对象")
    footer(s)
    steps = [
        ("输入", "关键词\n源站列表\n平台选择", BLUE_FILL, BLUE),
        ("趋势采集", "CollectRequest\n→ TrendItem[]", TEAL_FILL, TEAL),
        ("内容生成", "GenerateRequest\n→ PostDraft", GOLD_FILL, GOLD),
        ("发布预览", "PreviewRequest\n→ AdapterPreview", PURPLE_FILL, PURPLE),
        ("审计记录", "JobEvent\nAuditMetric", RED_FILL, RED),
        ("实验产物", "run.json\nreport.md\nscreenshots", FILL, NAVY),
    ]
    x = 0.55
    for i, (head, body, fill, line) in enumerate(steps):
        rect(s, x, 2.0, 1.75, 0.45, head, fill, line, line, size=13, bold=True)
        rect(s, x, 2.52, 1.75, 1.05, body, RGBColor(255, 255, 255), line, INK, size=11)
        if i < len(steps) - 1:
            arrow(s, x + 1.76, 2.78, x + 2.18, 2.78, BLUE)
        x += 2.05
    label(s, 0.85, 4.45, 11.4, 0.8, "输出流要点：接口返回结构化对象，页面展示用户可读状态，运行日志和未来 AuditMetric 支撑课程实验复盘。", NAVY, size=14, align=PP_ALIGN.CENTER)
    save(prs, "02_output_flow")


def api_sequence():
    prs, s = new_prs()
    title(s, "核心接口调用时序", "前端按采集 → 生成 → 预览顺序调用后端，后端服务写入内存仓库并返回 JobEvent")
    footer(s)
    actors = [
        ("运营人员", 0.7),
        ("React 前端", 2.7),
        ("FastAPI API", 4.9),
        ("Service 层", 7.1),
        ("Repository", 9.3),
        ("Adapter", 11.1),
    ]
    for name, x in actors:
        rect(s, x, 1.45, 1.45, 0.48, name, BLUE_FILL, BLUE, BLUE, size=11.5, bold=True)
        line = s.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(x + 0.72), Inches(2.0), Inches(x + 0.72), Inches(6.25))
        line.line.color.rgb = BORDER
        line.line.width = Pt(1)
    calls = [
        (1, 2, 2.25, "POST /api/trends/collect"),
        (2, 3, 2.75, "TrendService.collect"),
        (3, 4, 3.22, "replace_trends + add_job"),
        (1, 2, 3.78, "POST /api/content/generate"),
        (2, 3, 4.28, "ContentService.generate"),
        (3, 4, 4.75, "add_post + add_job"),
        (1, 2, 5.30, "POST /api/publish/preview"),
        (2, 5, 5.80, "adapter.preview(post)"),
    ]
    positions = [x + 0.72 for _, x in actors]
    for a, b, y, text in calls:
        arrow(s, positions[a], y, positions[b], y, TEAL if y < 3.5 else PURPLE)
        label(s, min(positions[a], positions[b]) + 0.08, y - 0.25, abs(positions[b] - positions[a]) - 0.15, 0.22, text, INK, size=8.8, align=PP_ALIGN.CENTER)
    save(prs, "03_api_sequence")


def data_model():
    prs, s = new_prs()
    title(s, "核心数据结构关系", "当前使用 Pydantic 模型和内存仓库，后续可映射到 SQLite / Postgres")
    footer(s)
    rect(s, 0.75, 1.65, 2.25, 0.95, "TrendSource\nid/name/category/enabled", BLUE_FILL, BLUE, BLUE, size=11, bold=True)
    rect(s, 3.55, 1.65, 2.25, 0.95, "TrendItem\nid/rank/title/source/score", BLUE_FILL, BLUE, BLUE, size=11, bold=True)
    rect(s, 6.35, 1.65, 2.25, 0.95, "PostDraft\ntrend_id/status/variants", TEAL_FILL, TEAL, TEAL, size=11, bold=True)
    rect(s, 9.15, 1.65, 2.25, 0.95, "ContentVariant\nplatform/title/body/tags", TEAL_FILL, TEAL, TEAL, size=11, bold=True)
    rect(s, 2.15, 4.25, 2.35, 0.95, "AdapterStatus\nmode/state/account", PURPLE_FILL, PURPLE, PURPLE, size=11, bold=True)
    rect(s, 5.30, 4.25, 2.35, 0.95, "AdapterPreview\nmessage/command/payload", PURPLE_FILL, PURPLE, PURPLE, size=11, bold=True)
    rect(s, 8.45, 4.25, 2.35, 0.95, "JobEvent\nkind/title/status/time", GOLD_FILL, GOLD, GOLD, size=11, bold=True)
    arrow(s, 3.0, 2.12, 3.55, 2.12)
    arrow(s, 5.8, 2.12, 6.35, 2.12)
    arrow(s, 8.6, 2.12, 9.15, 2.12)
    arrow(s, 7.48, 2.60, 6.45, 4.25, TEAL)
    arrow(s, 4.5, 4.72, 5.30, 4.72, PURPLE)
    arrow(s, 7.65, 4.72, 8.45, 4.72, GOLD)
    label(s, 1.0, 6.15, 11.2, 0.36, "关系说明：TrendItem 是内容生成输入；PostDraft 聚合多平台 ContentVariant；AdapterPreview 是发布管理页面展示和未来平台接入的边界对象。", NAVY, size=12.5, align=PP_ALIGN.CENTER)
    save(prs, "04_data_model")


def deployment():
    prs, s = new_prs()
    title(s, "运行部署拓扑", "当前为本地课程实验部署，后续可替换持久化存储与正式平台账号")
    footer(s)
    rect(s, 0.9, 1.65, 2.0, 0.8, "浏览器\nhttp://127.0.0.1:5173", BLUE_FILL, BLUE, BLUE, size=11, bold=True)
    rect(s, 3.55, 1.65, 2.2, 0.8, "Vite Dev Server\nReact 控制台", BLUE_FILL, BLUE, BLUE, size=11, bold=True)
    rect(s, 6.35, 1.65, 2.2, 0.8, "Uvicorn/FastAPI\nhttp://127.0.0.1:8000", TEAL_FILL, TEAL, TEAL, size=11, bold=True)
    rect(s, 9.25, 1.05, 2.55, 0.68, "external/wemp-operator", GOLD_FILL, GOLD, GOLD, size=10.5, bold=True)
    rect(s, 9.25, 1.95, 2.55, 0.68, "external/toutiao", GOLD_FILL, GOLD, GOLD, size=10.5, bold=True)
    rect(s, 9.25, 2.85, 2.55, 0.68, "external/XiaohongshuSkills", GOLD_FILL, GOLD, GOLD, size=10.5, bold=True)
    rect(s, 3.6, 4.4, 2.2, 0.75, "MemoryRepository\n当前内存状态", RED_FILL, RED, RED, size=11, bold=True)
    rect(s, 6.45, 4.4, 2.2, 0.75, "未来 SQLite/Postgres\n持久化任务与草稿", FILL, NAVY, NAVY, size=10.5, bold=True)
    rect(s, 9.25, 4.4, 2.55, 0.75, "artifacts/runs\n报告、截图、审计产物", FILL, NAVY, NAVY, size=10.5, bold=True)
    arrow(s, 2.9, 2.05, 3.55, 2.05)
    arrow(s, 5.75, 2.05, 6.35, 2.05)
    arrow(s, 8.55, 2.05, 9.25, 1.39, TEAL)
    arrow(s, 8.55, 2.05, 9.25, 2.29, TEAL)
    arrow(s, 8.55, 2.05, 9.25, 3.19, TEAL)
    arrow(s, 7.45, 2.45, 4.7, 4.4, RED)
    arrow(s, 5.8, 4.78, 6.45, 4.78, NAVY)
    arrow(s, 8.65, 4.78, 9.25, 4.78, NAVY)
    save(prs, "05_deployment")


def module_io():
    prs, s = new_prs()
    title(s, "模块输入输出设计", "每个模块都明确输入、处理、输出，方便审计和后续替换")
    footer(s)
    headers = ["模块", "输入", "处理", "输出"]
    xs = [0.55, 2.9, 5.6, 8.55]
    ws = [2.1, 2.4, 2.65, 3.9]
    for x, w, h in zip(xs, ws, headers):
        rect(s, x, 1.35, w, 0.42, h, BLUE_FILL, BLUE, BLUE, size=11.5, bold=True)
    rows = [
        ("趋势洞察", "sources / keyword / limit", "外部爬虫 + 规范化 + fallback", "TrendItem[] + CollectResponse"),
        ("内容工作室", "trend_id / platforms", "模板生成三平台版本", "PostDraft + ContentVariant[]"),
        ("发布管理", "post_id / platform", "AdapterRegistry 分发 preview", "AdapterPreview + JobEvent"),
        ("运行日志", "服务事件", "add_job 截断保留最近 80 条", "JobEvent[]"),
    ]
    y = 1.88
    fills = [TEAL_FILL, GOLD_FILL, PURPLE_FILL, FILL]
    lines = [TEAL, GOLD, PURPLE, NAVY]
    for idx, row in enumerate(rows):
        for col, text in enumerate(row):
            rect(s, xs[col], y, ws[col], 0.72, text, fills[idx], lines[idx], INK if col else lines[idx], size=10.3, bold=(col == 0))
        y += 0.9
    rect(s, 0.85, 5.9, 11.45, 0.62, "实现原则：当前模块先保证结构化预览和可演示流程；后续通过审计指标、持久化和审核状态逐步升级到可发布闭环。", RED_FILL, RED, RED, size=12.3, bold=True)
    save(prs, "06_module_io")


def agent_skill_architecture():
    prs, s = new_prs()
    title(s, "Agent 与 Skill 架构图", "前后端只是承载层，核心是任务 Agent 调度多个平台 Skill 完成信息获取、内容生成与发布预览")
    footer(s)

    rect(s, 0.55, 1.35, 1.7, 0.58, "任务目标\n关键词/平台/边界", BLUE_FILL, BLUE, BLUE, size=10.5, bold=True)
    rect(s, 2.75, 1.25, 2.05, 0.78, "CoordinatorAgent\n任务拆解与状态机", TEAL_FILL, TEAL, TEAL, size=11, bold=True)
    rect(s, 5.35, 1.25, 2.2, 0.78, "Agent Orchestrator\n计划/调用/审计", TEAL_FILL, TEAL, TEAL, size=11, bold=True)
    rect(s, 8.1, 1.25, 2.05, 0.78, "SkillRegistry\n能力发现与路由", GOLD_FILL, GOLD, GOLD, size=11, bold=True)
    rect(s, 10.65, 1.35, 1.9, 0.58, "前端/后端\n展示与 API", BLUE_FILL, BLUE, BLUE, size=10.5, bold=True)

    agents = [
        ("TrendScoutAgent\n热点采集", 0.72, 3.0, TEAL_FILL, TEAL),
        ("InsightAnalystAgent\n洞察评分", 2.78, 3.0, TEAL_FILL, TEAL),
        ("BriefPlannerAgent\n选题策略", 4.84, 3.0, GOLD_FILL, GOLD),
        ("PlatformWriterAgent\n平台改写", 6.90, 3.0, GOLD_FILL, GOLD),
        ("ComplianceReviewer\n审核门禁", 8.96, 3.0, RED_FILL, RED),
        ("PublisherAgent\n预览发布", 11.02, 3.0, PURPLE_FILL, PURPLE),
    ]
    for idx, (text, x, y, fill, line) in enumerate(agents):
        rect(s, x, y, 1.55, 0.72, text, fill, line, line, size=9.8, bold=True)
        if idx < len(agents) - 1:
            arrow(s, x + 1.55, y + 0.36, x + 2.06, y + 0.36, BLUE)

    skills = [
        ("wemp-operator\nsmart-collect", 1.0, 5.05),
        ("ToutiaoSkill\nscraper/preview", 3.35, 5.05),
        ("XiaohongshuSkill\nsearch/preview", 5.70, 5.05),
        ("WechatSkill\nmarkdown/draft", 8.05, 5.05),
        ("ImageSkill\nimage_prompt/素材", 10.40, 5.05),
    ]
    for text, x, y in skills:
        rect(s, x, y, 1.9, 0.68, text, FILL, NAVY, NAVY, size=9.5, bold=True)

    arrow(s, 2.25, 1.64, 2.75, 1.64)
    arrow(s, 4.8, 1.64, 5.35, 1.64, TEAL)
    arrow(s, 7.55, 1.64, 8.1, 1.64, GOLD)
    arrow(s, 10.15, 1.64, 10.65, 1.64, BLUE)
    rect(s, 0.82, 6.35, 11.45, 0.48, "审计对象：AgentTask、SkillCall、InsightCard、ContentBrief、ContentVariant、ReviewResult、AdapterPreview、JobEvent", RED_FILL, RED, RED, size=10.8, bold=True)
    save(prs, "07_agent_skill_architecture")

    im = Image.new("RGB", (PNG_W, PNG_H), "white")
    d = ImageDraw.Draw(im)
    png_title(d, "Agent 与 Skill 架构图", "任务由 Agent 编排，平台能力通过 SkillRegistry 统一接入，前后端只负责展示和 API 承载")
    blue = pil_color((36, 99, 170))
    teal = pil_color((15, 118, 110))
    gold = pil_color((180, 120, 30))
    purple = pil_color((103, 80, 164))
    red = pil_color((155, 28, 28))
    ink = pil_color((31, 41, 55))
    border = pil_color((203, 213, 225))
    png_rect(d, (90, 235, 345, 330), "任务目标\n关键词 / 平台 / 边界", (232, 241, 252), blue, blue, 25, True)
    png_rect(d, (430, 220, 730, 345), "CoordinatorAgent\n任务拆解与状态机", (228, 247, 244), teal, teal, 25, True)
    png_rect(d, (820, 220, 1145, 345), "Agent Orchestrator\n计划 / 调用 / 审计", (228, 247, 244), teal, teal, 25, True)
    png_rect(d, (1235, 220, 1535, 345), "SkillRegistry\n能力发现与路由", (255, 247, 225), gold, gold, 25, True)
    png_rect(d, (1620, 235, 1840, 330), "前端 / 后端\n展示与 API", (232, 241, 252), blue, blue, 24, True)
    for start, end, color in [
        ((345, 282), (430, 282), blue),
        ((730, 282), (820, 282), teal),
        ((1145, 282), (1235, 282), gold),
        ((1535, 282), (1620, 282), blue),
    ]:
        png_arrow(d, start, end, color)

    agent_boxes = [
        ((90, 455, 325, 560), "TrendScoutAgent\n热点采集", teal, (228, 247, 244)),
        ((390, 455, 625, 560), "InsightAnalystAgent\n洞察评分", teal, (228, 247, 244)),
        ((690, 455, 925, 560), "BriefPlannerAgent\n选题策略", gold, (255, 247, 225)),
        ((990, 455, 1225, 560), "PlatformWriterAgent\n平台改写", gold, (255, 247, 225)),
        ((1290, 455, 1525, 560), "ComplianceReviewer\n审核门禁", red, (255, 241, 242)),
        ((1590, 455, 1825, 560), "PublisherAgent\n预览发布", purple, (243, 238, 255)),
    ]
    for idx, (box, text, line, fill) in enumerate(agent_boxes):
        png_rect(d, box, text, fill, line, line, 23, True)
        if idx < len(agent_boxes) - 1:
            png_arrow(d, (box[2], 508), (agent_boxes[idx + 1][0][0], 508), blue, 4)

    skill_boxes = [
        ((140, 715, 405, 820), "wemp-operator\nsmart-collect", ink),
        ((470, 715, 735, 820), "ToutiaoSkill\nscraper / preview", ink),
        ((800, 715, 1065, 820), "XiaohongshuSkill\nsearch / preview", ink),
        ((1130, 715, 1395, 820), "WechatSkill\nmarkdown / draft", ink),
        ((1460, 715, 1725, 820), "ImageSkill\nprompt / material", ink),
    ]
    for box, text, color in skill_boxes:
        png_rect(d, box, text, (248, 250, 252), border, color, 23, True)
    png_rect(
        d,
        (170, 905, 1750, 975),
        "审计对象：AgentTask、SkillCall、InsightCard、ContentBrief、ContentVariant、ReviewResult、AdapterPreview、JobEvent",
        (255, 241, 242),
        red,
        red,
        25,
        True,
        radius=16,
    )
    im.save(OUT / "07_agent_skill_architecture.pptx.png")


def end_to_end_agent_flow():
    prs, s = new_prs()
    title(s, "获取信息、处理信息、发布信息全流程", "每个阶段都明确 Agent、Skill、输入输出和审计证据")
    footer(s)
    cols = [
        ("1 获取信息", "TrendScoutAgent\n调用采集 Skill", "输入：关键词、源站、limit\n输出：RawTrend / TrendItem\nEvidenceItem", TEAL_FILL, TEAL),
        ("2 处理信息", "InsightAnalyst + BriefPlanner\n洞察、选题、证据整理", "输入：TrendItem、历史记录\n平台范围\n输出：InsightCard / ContentBrief", BLUE_FILL, BLUE),
        ("3 生成内容", "PlatformWriter + MediaPrompt\n按平台画像改写", "输入：ContentBrief\nPlatformProfile\n输出：ContentVariant\nImagePromptSpec", GOLD_FILL, GOLD),
        ("4 发布信息", "ComplianceReviewer\nPublisherAgent\n审核、预览、草稿边界", "输入：ContentVariant\nReviewResult\n输出：AdapterPreview\nJobEvent", PURPLE_FILL, PURPLE),
    ]
    x = 0.62
    for idx, (head, agent, io, fill, line) in enumerate(cols):
        rect(s, x, 1.45, 2.65, 0.45, head, fill, line, line, size=12.5, bold=True)
        rect(s, x, 2.05, 2.65, 0.80, agent, RGBColor(255, 255, 255), line, INK, size=10.8, bold=True)
        rect(s, x, 3.08, 2.65, 1.55, io, FILL, line, INK, size=9.7)
        rect(s, x, 4.92, 2.65, 0.58, "审计：SkillCall + Trace", RED_FILL, RED, RED, size=9.6, bold=True)
        if idx < len(cols) - 1:
            arrow(s, x + 2.65, 3.24, x + 3.00, 3.24, BLUE)
        x += 3.05
    label(s, 0.88, 6.12, 11.3, 0.42, "关键原则：没有 Evidence 不生成观点；没有 Brief 不直接写稿；没有 ReviewResult 不进入平台预览；没有 approved 不允许真实发布。", NAVY, size=12.5, align=PP_ALIGN.CENTER)
    save(prs, "08_end_to_end_agent_flow")

    im = Image.new("RGB", (PNG_W, PNG_H), "white")
    d = ImageDraw.Draw(im)
    png_title(d, "获取信息、处理信息、发布信息全流程", "每个阶段都明确 Agent、Skill、输入输出和审计证据")
    colors = [
        ((228, 247, 244), pil_color((15, 118, 110))),
        ((232, 241, 252), pil_color((36, 99, 170))),
        ((255, 247, 225), pil_color((180, 120, 30))),
        ((243, 238, 255), pil_color((103, 80, 164))),
    ]
    blocks = [
        ("1 获取信息", "TrendScoutAgent\n调用采集 Skill", "输入：关键词、源站、limit\n输出：RawTrend / TrendItem\nEvidenceItem"),
        ("2 处理信息", "InsightAnalyst + BriefPlanner\n洞察、选题、证据整理", "输入：TrendItem、历史记录\n平台范围\n输出：InsightCard / ContentBrief"),
        ("3 生成内容", "PlatformWriter + MediaPrompt\n按平台画像改写", "输入：ContentBrief\nPlatformProfile\n输出：ContentVariant\nImagePromptSpec"),
        ("4 发布信息", "ComplianceReviewer\nPublisherAgent\n审核、预览、草稿边界", "输入：ContentVariant\nReviewResult\n输出：AdapterPreview\nJobEvent"),
    ]
    x = 95
    for idx, (head, agent, io) in enumerate(blocks):
        fill, line = colors[idx]
        png_rect(d, (x, 245, x + 365, 315), head, fill, line, line, 25, True, radius=16)
        png_rect(d, (x, 350, x + 365, 465), agent, (255, 255, 255), line, pil_color((31, 41, 55)), 23, True, radius=16)
        png_rect(d, (x, 505, x + 365, 690), io, (248, 250, 252), line, pil_color((31, 41, 55)), 21, False, radius=16)
        png_rect(d, (x, 735, x + 365, 810), "审计：SkillCall + Trace", (255, 241, 242), pil_color((155, 28, 28)), pil_color((155, 28, 28)), 22, True, radius=16)
        if idx < len(blocks) - 1:
            png_arrow(d, (x + 365, 575), (x + 445, 575), pil_color((36, 99, 170)), 5)
        x += 445
    png_rect(
        d,
        (155, 910, 1765, 985),
        "原则：没有 Evidence 不生成观点；没有 Brief 不直接写稿；没有 ReviewResult 不进入平台预览；没有 approved 不允许真实发布。",
        (248, 250, 252),
        pil_color((15, 42, 67)),
        pil_color((15, 42, 67)),
        24,
        True,
        radius=18,
    )
    im.save(OUT / "08_end_to_end_agent_flow.pptx.png")


if __name__ == "__main__":
    system_architecture()
    output_flow()
    api_sequence()
    data_model()
    deployment()
    module_io()
    agent_skill_architecture()
    end_to_end_agent_flow()
    print(f"wrote diagrams to {PPTX_DIR}")
